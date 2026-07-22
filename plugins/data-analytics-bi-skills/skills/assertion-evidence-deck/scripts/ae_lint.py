#!/usr/bin/env python3
"""Audit a .pptx against the assertion-evidence checklist.

Usage:
  python ae_lint.py deck.pptx [--json] [--max-body-words N]

Reports headline violations, bullet characters, body word-count overruns, banned
formatting (italics, underline, shouting caps), missing source tags on slides
that carry figures, and unexpected fonts. Exits non-zero when it finds an ERROR
(warnings do not fail the run). Fix the generator and rebuild rather than
hand-patching the packed XML.

Requires: python-pptx  (pip install python-pptx)
"""
from __future__ import annotations

import argparse
import json
import re
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    sys.exit("python-pptx is required: pip install python-pptx")

HEAD_LINE_CHARS = 52          # ~105 chars over two lines at 28 pt / 11.52 in
MAX_BODY_WORDS = 40
BOTTOM_ZONE_IN = 6.5          # shapes below this top are source tags / slide numbers
ALLOWED_FONTS = {"calibri", "arial", "gotham", "montserrat", "goudy old style"}
BULLET_CHARS = "•‣▪◦·●■–—"
BULLET_RE = re.compile(r"^\s*[" + re.escape(BULLET_CHARS) + r"]\s+")
DASH_BULLET_RE = re.compile(r"^\s*[-*]\s+")
DIGIT_RE = re.compile(r"\d")
# All-caps words worth flagging (real words shouted), sparing short acronyms.
CAPS_RE = re.compile(r"\b[A-Z]{6,}\b")
CAPS_ALLOW = {"AUTORECON"}    # domain terms that are legitimately upper


def estimate_lines(text, max_chars=HEAD_LINE_CHARS):
    lines, cur = 0, ""
    for word in str(text).split():
        cand = (cur + " " + word).strip()
        if len(cand) > max_chars and cur:
            lines += 1
            cur = word
        else:
            cur = cand
    if cur:
        lines += 1
    return max(lines, 1)


def _emu_to_in(v):
    return (v or 0) / 914400.0


class Findings:
    def __init__(self):
        self.items = []

    def add(self, slide, level, check, msg):
        self.items.append({"slide": slide, "level": level, "check": check, "message": msg})

    @property
    def errors(self):
        return [f for f in self.items if f["level"] == "error"]


def _text_shapes(slide):
    """Yield (shape, top_in) for shapes with a text frame, non-empty."""
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            yield sh, _emu_to_in(sh.top)


def lint_slide(idx, slide, fnd, max_words):
    texts = list(_text_shapes(slide))
    # Headline = highest (smallest top) text shape not in the bottom zone.
    candidates = [(sh, t) for sh, t in texts if t < BOTTOM_ZONE_IN]
    headline_shape = min(candidates, key=lambda p: p[1])[0] if candidates else None
    headline = headline_shape.text_frame.text.strip() if headline_shape else ""

    if not headline:
        fnd.add(idx, "error", "headline", "no headline text found on the slide")
    else:
        lines = estimate_lines(headline)
        if lines > 2 or len(headline) > 110:
            fnd.add(idx, "error", "headline",
                    f"headline exceeds two lines (~{lines}): {headline!r}")
        if len(headline.split()) < 3:
            fnd.add(idx, "warn", "headline",
                    f"headline looks like a topic label, not a sentence: {headline!r}")

    body_words = 0
    has_source = False
    has_figure = False

    for sh, top in texts:
        tf = sh.text_frame
        is_bottom = top >= BOTTOM_ZONE_IN
        txt = tf.text.strip()
        if txt.lower().startswith("source:"):
            has_source = True
        # per-paragraph and per-run checks
        for para in tf.paragraphs:
            ptext = "".join(r.text for r in para.runs) or para.text
            if BULLET_RE.match(ptext) or (DASH_BULLET_RE.match(ptext) and sh is not headline_shape):
                fnd.add(idx, "error", "bullet",
                        f"bulleted line under a headline: {ptext.strip()[:60]!r}")
            # bullet formatting carried in XML (buChar / buAutoNum)
            pPr = para._pPr
            if pPr is not None and (pPr.find(_qn("a:buChar")) is not None or
                                    pPr.find(_qn("a:buAutoNum")) is not None):
                fnd.add(idx, "warn", "bullet",
                        "paragraph carries a bullet/auto-number format; use spatial layout instead")
            for run in para.runs:
                if sh is headline_shape or is_bottom:
                    continue
                if run.font.italic:
                    fnd.add(idx, "error", "format", f"italic text: {run.text.strip()[:40]!r}")
                if run.font.underline:
                    fnd.add(idx, "error", "format", f"underlined text: {run.text.strip()[:40]!r}")
                name = run.font.name
                if name and name.lower() not in ALLOWED_FONTS:
                    fnd.add(idx, "warn", "font",
                            f"font {name!r} is not in the sanctioned set "
                            f"({', '.join(sorted(ALLOWED_FONTS))})")
            for m in CAPS_RE.findall(ptext):
                if m not in CAPS_ALLOW:
                    fnd.add(idx, "warn", "caps", f"all-caps word: {m!r} (avoid all capitals)")
        if sh is not headline_shape and not is_bottom:
            body_words += len(txt.split())
        if DIGIT_RE.search(txt) and sh is not headline_shape and not is_bottom:
            has_figure = True

    # tables and charts carry figures too
    for sh in slide.shapes:
        if getattr(sh, "has_table", False) and sh.has_table:
            for row in sh.table.rows:
                for cell in row.cells:
                    if DIGIT_RE.search(cell.text):
                        has_figure = True
        if getattr(sh, "has_chart", False) and sh.has_chart:
            has_figure = True

    if body_words > max_words:
        fnd.add(idx, "error", "wordcount",
                f"{body_words} body words (max {max_words}); the claim is too big — split it")
    if has_figure and not has_source:
        fnd.add(idx, "warn", "source",
                "slide carries a figure but has no 'Source:' tag")


def _qn(tag):
    from pptx.oxml.ns import qn
    return qn(tag)


def lint(path, max_words=MAX_BODY_WORDS):
    prs = Presentation(path)
    fnd = Findings()
    for i, slide in enumerate(prs.slides, start=1):
        lint_slide(i, slide, fnd, max_words)
    return fnd


def main(argv=None):
    ap = argparse.ArgumentParser(description="Lint a .pptx against the assertion-evidence checklist.")
    ap.add_argument("pptx", help="path to the .pptx to audit")
    ap.add_argument("--json", action="store_true", help="machine-readable output")
    ap.add_argument("--max-body-words", type=int, default=MAX_BODY_WORDS)
    args = ap.parse_args(argv)

    fnd = lint(args.pptx, args.max_body_words)

    if args.json:
        print(json.dumps({"findings": fnd.items,
                          "errors": len(fnd.errors),
                          "total": len(fnd.items)}, indent=2))
    else:
        if not fnd.items:
            print("clean — no assertion-evidence violations found")
        for f in fnd.items:
            print(f"[{f['level'].upper():5}] slide {f['slide']:>2} {f['check']:<9} {f['message']}")
        print(f"\n{len(fnd.errors)} error(s), {len(fnd.items) - len(fnd.errors)} warning(s)")

    return 1 if fnd.errors else 0


if __name__ == "__main__":
    sys.exit(main())
