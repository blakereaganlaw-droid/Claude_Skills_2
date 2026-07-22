#!/usr/bin/env python3
"""Build a compliant assertion-evidence .pptx from a JSON deck spec.

Usage:
  python build_deck.py deck_spec.json -o output.pptx [--brand ut|neutral] [--font NAME]
  python build_deck.py --schema         # print the spec format and the eight slide kinds

The builder encodes the geometry, typography, and colors verified in
references/design-tokens.md, so slides come out compliant without measuring by
hand. It refuses a headline that busts the two-line budget rather than letting it
overflow silently, and it warns when a chosen text/background pair falls below
WCAG AA contrast.

Requires: python-pptx  (pip install python-pptx)
"""
from __future__ import annotations

import argparse
import json
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
except ImportError:
    sys.exit("python-pptx is required: pip install python-pptx")

# --- Geometry (inches), verified from the official template ---------------
CANVAS_W, CANVAS_H = 13.333, 7.5
MARGIN = 0.92
HEAD = dict(x=0.92, y=0.62, w=11.52, h=0.94)      # 28 pt bold, left
BODY = dict(x=0.92, y=1.95, w=11.52, h=4.55)      # 18 pt
COL = dict(y=1.95, w=3.53, h=4.30, left_x=2.15, right_x=7.65)
SRC = dict(x=0.92, y=6.95, w=8.0, h=0.40)         # 12 pt muted, left
NUM = dict(x=9.42, y=6.95, w=3.0, h=0.40)         # 12 pt muted, right
TITLE_HEAD = dict(x=0.0, y=1.67, w=13.333, h=2.22)  # 40 pt

HEAD_PT, BODY_PT, SRC_PT, TITLE_PT = 28, 18, 12, 40
HEAD_LINE_CHARS = 52          # ~105 chars over two lines at 28 pt / 11.52 in
BANNED = {"kind"}             # reserved keys, not content

# --- Palettes -------------------------------------------------------------
PALETTES = {
    "ut": dict(primary="FF8200", ink="4B4B4B", muted="767676",
               secondary="A6A6A6", bg="FFFFFF"),
    "neutral": dict(primary="4B4B4B", ink="333333", muted="767676",
                    secondary="A6A6A6", bg="FFFFFF"),
}

SLIDE_KINDS = {
    "title":      "Title slide. Fields: headline, subtitle.",
    "statement":  "Headline + short body text. Fields: headline, body, source, notes.",
    "magnitude":  "2-4 large numbers. Fields: headline, numbers:[{value,label}], source, notes.",
    "chart":      "Native chart. Fields: headline, chart:{type:bar|column|line|stacked, categories:[], series:[{name,values:[]}]}, source, notes.",
    "table":      "Table. Fields: headline, columns:[], rows:[[]], emphasize_column:int, source, notes.",
    "two_column": "Defect pair. Fields: headline, left:{title,items:[]}, right:{title,items:[]}, source, notes.",
    "flow":       "Left-to-right process. Fields: headline, steps:[], source, notes.",
    "summary":    "Restate the main claim. Fields: headline, points:[], notes.",
}


def rgb(hex6):
    return RGBColor.from_string(hex6)


# --- WCAG contrast --------------------------------------------------------
def _lin(c):
    c /= 255.0
    return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4


def _lum(hex6):
    r, g, b = (int(hex6[i:i + 2], 16) for i in (0, 2, 4))
    return 0.2126 * _lin(r) + 0.7152 * _lin(g) + 0.0722 * _lin(b)


def contrast(fg, bg):
    a, b = _lum(fg), _lum(bg)
    hi, lo = max(a, b), min(a, b)
    return (hi + 0.05) / (lo + 0.05)


# --- Headline budget ------------------------------------------------------
def estimate_lines(text, max_chars=HEAD_LINE_CHARS):
    lines, cur = 0, ""
    for word in str(text).split():
        cand = (cur + " " + word).strip()
        if len(cand) > max_chars and cur:
            lines += 1
            cur = word
        else:
            cur = cand
    if cur:
        lines += 1
    return max(lines, 1)


# --- Low-level helpers ----------------------------------------------------
def _textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def _style(run, font, size, color, bold=False):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = False
    run.font.underline = False
    run.font.color.rgb = rgb(color)


def _set_para(tf, text, font, size, color, bold=False, align=PP_ALIGN.LEFT, first=True):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.line_spacing = 0.9
    run = p.add_run()
    run.text = str(text)
    _style(run, font, size, color, bold)
    return p


def add_headline(slide, text, pal, font, title_slide=False):
    g = TITLE_HEAD if title_slide else HEAD
    size = TITLE_PT if title_slide else HEAD_PT
    _, tf = _textbox(slide, g["x"], g["y"], g["w"], g["h"], MSO_ANCHOR.BOTTOM)
    _set_para(tf, text, font, size, pal["ink"], bold=True,
              align=PP_ALIGN.CENTER if title_slide else PP_ALIGN.LEFT)


def add_source(slide, text, pal, font):
    if not text:
        return
    _, tf = _textbox(slide, SRC["x"], SRC["y"], SRC["w"], SRC["h"])
    _set_para(tf, text, font, SRC_PT, pal["muted"])


def add_number(slide, n, pal, font):
    _, tf = _textbox(slide, NUM["x"], NUM["y"], NUM["w"], NUM["h"])
    _set_para(tf, str(n), font, SRC_PT, pal["muted"], align=PP_ALIGN.RIGHT)


def add_notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = str(text)


# --- Slide builders -------------------------------------------------------
def build_title(slide, spec, pal, font):
    add_headline(slide, spec["headline"], pal, font, title_slide=True)
    if spec.get("subtitle"):
        _, tf = _textbox(slide, 0.0, 4.0, CANVAS_W, 1.0)
        _set_para(tf, spec["subtitle"], font, 22, pal["muted"], align=PP_ALIGN.CENTER)


def build_statement(slide, spec, pal, font):
    _, tf = _textbox(slide, BODY["x"], BODY["y"], BODY["w"], BODY["h"])
    first = True
    for para in str(spec.get("body", "")).split("\n"):
        _set_para(tf, para, font, BODY_PT, pal["ink"], first=first)
        first = False


def build_magnitude(slide, spec, pal, font):
    nums = spec.get("numbers", [])[:4]
    if not nums:
        return
    n = len(nums)
    gap = 0.4
    total = BODY["w"]
    cell = (total - gap * (n - 1)) / n
    for i, item in enumerate(nums):
        x = BODY["x"] + i * (cell + gap)
        _, tf = _textbox(slide, x, 2.6, cell, 2.0, MSO_ANCHOR.MIDDLE)
        _set_para(tf, item.get("value", ""), font, 54, pal["primary"], bold=True,
                  align=PP_ALIGN.CENTER)
        _set_para(tf, item.get("label", ""), font, BODY_PT, pal["ink"],
                  align=PP_ALIGN.CENTER, first=False)


_CHART_TYPES = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "stacked": XL_CHART_TYPE.COLUMN_STACKED,
}


def build_chart(slide, spec, pal, font):
    c = spec.get("chart", {})
    cd = CategoryChartData()
    cd.categories = c.get("categories", [])
    for s in c.get("series", []):
        cd.add_series(s.get("name", ""), s.get("values", []))
    ctype = _CHART_TYPES.get(c.get("type", "column"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    gframe = slide.shapes.add_chart(
        ctype, Inches(BODY["x"]), Inches(BODY["y"]),
        Inches(BODY["w"]), Inches(BODY["h"] - 0.3), cd)
    chart = gframe.chart
    chart.has_legend = len(c.get("series", [])) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    palette = [pal["primary"], pal["secondary"], pal["muted"]]
    for i, series in enumerate(chart.series):
        try:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = rgb(palette[i % len(palette)])
        except Exception:
            pass
    try:
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(12)
        plot.data_labels.font.name = font
        plot.data_labels.font.color.rgb = rgb(pal["ink"])
    except Exception:
        pass


def build_table(slide, spec, pal, font):
    cols = spec.get("columns", [])
    rows = spec.get("rows", [])
    emph = spec.get("emphasize_column")
    nrows, ncols = len(rows) + 1, max(1, len(cols))
    gtable = slide.shapes.add_table(
        nrows, ncols, Inches(BODY["x"]), Inches(BODY["y"]),
        Inches(BODY["w"]), Inches(min(BODY["h"], 0.5 * nrows + 0.3))).table
    for j, col in enumerate(cols):
        cell = gtable.cell(0, j)
        cell.text = str(col)
        _style(cell.text_frame.paragraphs[0].add_run() if False else
               cell.text_frame.paragraphs[0].runs[0], font, 14, pal["bg"], bold=True)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(pal["ink"])
    for i, row in enumerate(rows, start=1):
        for j in range(ncols):
            cell = gtable.cell(i, j)
            cell.text = str(row[j]) if j < len(row) else ""
            run = cell.text_frame.paragraphs[0].runs[0] if cell.text_frame.paragraphs[0].runs \
                else cell.text_frame.paragraphs[0].add_run()
            _style(run, font, 14, pal["ink"], bold=(emph == j))


def _column(slide, x, spec_col, pal, font):
    _, tf = _textbox(slide, x, COL["y"], COL["w"], COL["h"])
    _set_para(tf, spec_col.get("title", ""), font, 20, pal["primary"], bold=True)
    for item in spec_col.get("items", []):
        _set_para(tf, item, font, 16, pal["ink"], first=False)


def build_two_column(slide, spec, pal, font):
    _column(slide, COL["left_x"] - COL["w"] / 2, spec.get("left", {}), pal, font)
    _column(slide, COL["right_x"] - COL["w"] / 2, spec.get("right", {}), pal, font)


def build_flow(slide, spec, pal, font):
    steps = spec.get("steps", [])[:7]
    if not steps:
        return
    n = len(steps)
    arrow = 0.45
    total = BODY["w"]
    box = (total - arrow * (n - 1)) / n
    y, h = 3.1, 1.3
    for i, step in enumerate(steps):
        x = BODY["x"] + i * (box + arrow)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(x), Inches(y), Inches(box), Inches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(pal["primary"] if i == n - 1 else pal["bg"])
        shp.line.color.rgb = rgb(pal["ink"])
        tf = shp.text_frame
        tf.word_wrap = True
        _set_para(tf, step, font, 14,
                  pal["ink"] if i != n - 1 else pal["ink"], bold=True,
                  align=PP_ALIGN.CENTER)
        if i < n - 1:
            _, atf = _textbox(slide, x + box, y, arrow, h, MSO_ANCHOR.MIDDLE)
            _set_para(atf, "→", font, 24, pal["muted"], align=PP_ALIGN.CENTER)


def build_summary(slide, spec, pal, font):
    _, tf = _textbox(slide, BODY["x"], BODY["y"], BODY["w"], BODY["h"])
    first = True
    for pt in spec.get("points", []):
        _set_para(tf, pt, font, 20, pal["ink"], first=first)
        first = False


BUILDERS = {
    "title": build_title, "statement": build_statement, "magnitude": build_magnitude,
    "chart": build_chart, "table": build_table, "two_column": build_two_column,
    "flow": build_flow, "summary": build_summary,
}


def build(spec, brand, font_override):
    pal = PALETTES[brand].copy()
    font = font_override or spec.get("font") or "Calibri"

    # Contrast guard on the standard text pairs.
    for label, fg in (("body/ink", pal["ink"]), ("source/muted", pal["muted"])):
        r = contrast(fg, pal["bg"])
        if r < 4.5:
            print(f"WARN: {label} {fg} on {pal['bg']} is {r:.2f}:1 (< 4.5:1 AA)",
                  file=sys.stderr)

    prs = Presentation()
    prs.slide_width = Inches(CANVAS_W)
    prs.slide_height = Inches(CANVAS_H)
    blank = prs.slide_layouts[6]

    slides = spec.get("slides", [])
    errors = []
    for i, s in enumerate(slides, start=1):
        kind = s.get("kind")
        if kind not in BUILDERS:
            errors.append(f"slide {i}: unknown kind {kind!r} (use one of {', '.join(SLIDE_KINDS)})")
            continue
        head = s.get("headline", "")
        if kind != "title" and not head:
            errors.append(f"slide {i} ({kind}): missing headline")
        if head and estimate_lines(head) > 2:
            errors.append(
                f"slide {i} ({kind}): headline exceeds the two-line budget "
                f"(~{estimate_lines(head)} lines). Split the slide or shorten: {head!r}")
    if errors:
        sys.exit("build refused:\n  - " + "\n  - ".join(errors))

    for i, s in enumerate(slides, start=1):
        kind = s["kind"]
        slide = prs.slides.add_slide(blank)
        if kind != "title":
            add_headline(slide, s.get("headline", ""), pal, font)
            add_source(slide, s.get("source"), pal, font)
            add_number(slide, i, pal, font)
        BUILDERS[kind](slide, s, pal, font)
        add_notes(slide, s.get("notes"))
    return prs


def print_schema():
    print("Deck spec (JSON):\n")
    print(json.dumps({
        "title": "string (used only if a title slide is present)",
        "font": "Calibri (optional; --font overrides)",
        "brand": "ut | neutral (optional; --brand overrides)",
        "slides": ["{kind: <one of below>, headline: ..., ...}"],
    }, indent=2))
    print("\nSlide kinds:\n")
    for k, v in SLIDE_KINDS.items():
        print(f"  {k:<11} {v}")
    print("\nEvery body slide also accepts: source (string), notes (speaker notes).")


def main(argv=None):
    ap = argparse.ArgumentParser(description="Build an assertion-evidence .pptx from a JSON spec.")
    ap.add_argument("spec", nargs="?", help="deck spec JSON file")
    ap.add_argument("-o", "--out", help="output .pptx path")
    ap.add_argument("--brand", choices=list(PALETTES), default="ut")
    ap.add_argument("--font", default=None, help="override the typeface (default Calibri)")
    ap.add_argument("--schema", action="store_true", help="print the spec format and exit")
    args = ap.parse_args(argv)

    if args.schema:
        print_schema()
        return
    if not args.spec or not args.out:
        ap.error("spec and -o/--out are required (or use --schema)")

    with open(args.spec, encoding="utf-8") as fh:
        spec = json.load(fh)
    prs = build(spec, args.brand, args.font)
    prs.save(args.out)
    print(f"wrote {args.out}  ({len(spec.get('slides', []))} slides, brand={args.brand})")


if __name__ == "__main__":
    main()
